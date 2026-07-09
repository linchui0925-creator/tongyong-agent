"""
文档解析服务，支持PDF/DOCX/DOC/XLSX/XLS/PPTX/PPT/TXT/MD/CSV等格式
"""
import csv
import io
import logging
from pathlib import Path
from typing import Optional, Tuple, List

logger = logging.getLogger(__name__)

# 依赖检查
_pdf_available = False
_docx_available = False
_xlsx_available = False
_pptx_available = False

try:
    from pypdf import PdfReader
    _pdf_available = True
except ImportError:
    logger.warning("pypdf not installed, PDF parsing disabled. Install with: pip install pypdf")

try:
    from docx import Document
    _docx_available = True
except ImportError:
    logger.warning("python-docx not installed, DOCX parsing disabled. Install with: pip install python-docx")

try:
    from openpyxl import load_workbook
    _xlsx_available = True
except ImportError:
    logger.warning("openpyxl not installed, Excel parsing disabled. Install with: pip install openpyxl")

try:
    from pptx import Presentation
    _pptx_available = True
except ImportError:
    logger.warning("python-pptx not installed, PPTX parsing disabled. Install with: pip install python-pptx")


def parse_pdf(path: Path) -> Tuple[str, dict]:
    """解析PDF文件，返回文本和元数据"""
    if not _pdf_available:
        return "", {"error": "pypdf not installed"}
    
    try:
        reader = PdfReader(str(path))
        text_parts = []
        meta = {
            "pages": len(reader.pages),
            "title": reader.metadata.title if reader.metadata else None,
            "author": reader.metadata.author if reader.metadata else None,
        }
        
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- 第{i+1}页 ---\n{page_text.strip()}")
        
        return "\n\n".join(text_parts), meta
    except Exception as e:
        logger.error(f"PDF解析失败: {e}")
        return "", {"error": str(e)}


def parse_docx(path: Path) -> Tuple[str, dict]:
    """解析DOCX Word文档"""
    if not _docx_available:
        return "", {"error": "python-docx not installed"}
    
    try:
        doc = Document(str(path))
        text_parts = []
        
        # 段落
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # 表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        meta = {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
        return "\n".join(text_parts), meta
    except Exception as e:
        logger.error(f"DOCX解析失败: {e}")
        return "", {"error": str(e)}


def parse_xlsx(path: Path, max_rows: int = 1000) -> Tuple[str, dict]:
    """解析Excel表格"""
    if not _xlsx_available:
        return "", {"error": "openpyxl not installed"}
    
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        text_parts = []
        meta = {"sheets": wb.sheetnames}
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"=== 工作表: {sheet_name} ===")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= max_rows:
                    text_parts.append(f"... (超过{max_rows}行，已截断)")
                    break
                row_text = " | ".join(str(cell) for cell in row if cell is not None and str(cell).strip())
                if row_text:
                    text_parts.append(row_text)
                row_count += 1
        
        return "\n".join(text_parts), meta
    except Exception as e:
        logger.error(f"Excel解析失败: {e}")
        return "", {"error": str(e)}


def parse_pptx(path: Path) -> Tuple[str, dict]:
    """解析PPTX演示文稿"""
    if not _pptx_available:
        return "", {"error": "python-pptx not installed"}
    
    try:
        prs = Presentation(str(path))
        text_parts = []
        meta = {"slides": len(prs.slides)}
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"--- 第{i+1}页幻灯片 ---\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_parts), meta
    except Exception as e:
        logger.error(f"PPTX解析失败: {e}")
        return "", {"error": str(e)}


def parse_text_file(path: Path) -> Tuple[str, dict]:
    """解析纯文本文件（txt/md/json/csv等）"""
    try:
        content = path.read_text(encoding="utf-8", errors="ignore")
        lines = content.splitlines()
        meta = {"lines": len(lines), "size": path.stat().st_size}
        return content, meta
    except Exception as e:
        logger.error(f"文本文件解析失败: {e}")
        return "", {"error": str(e)}


def parse_csv(path: Path, delimiter: str = ",") -> Tuple[str, dict]:
    """解析CSV文件"""
    try:
        text_parts = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f, delimiter=delimiter)
            for i, row in enumerate(reader):
                if i >= 1000:
                    text_parts.append("... (超过1000行，已截断)")
                    break
                row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                if row_text:
                    text_parts.append(row_text)
        return "\n".join(text_parts), {"rows": min(i+1, 1000)}
    except Exception as e:
        logger.error(f"CSV解析失败: {e}")
        return "", {"error": str(e)}
